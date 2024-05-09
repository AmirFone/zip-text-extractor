import zipfile
import os
import json
import io
import PyPDF2
from pptx import Presentation
from docx import Document
import pickle
from transformers import GPT2TokenizerFast

def extract_text_from_pkl(file):
    try:
        data = pickle.load(file)
        if not isinstance(data, str):
            data = str(data)
        return data
    except Exception as e:
        return f"Error occurred during pickle file processing: {e}"

def extract_text_from_pdf(file):
    pdf_reader = PyPDF2.PdfReader(file)
    text = ""
    for page in range(len(pdf_reader.pages)):
        text += pdf_reader.pages[page].extract_text()
    return text

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
    return text

def extract_text_from_docx(file):
    doc = Document(file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_zip(zip_file, file_path):
    result = {}
    for file_info in zip_file.infolist():
        if file_info.filename.startswith(file_path):
            relative_path = file_info.filename[len(file_path):]
            if relative_path and not relative_path.endswith('/'):
                with zip_file.open(file_info) as file:
                    try:
                        file_content = ""
                        if file_info.filename.lower().endswith('.pdf'):
                            file_content = extract_text_from_pdf(io.BytesIO(file.read()))
                        elif file_info.filename.lower().endswith('.pptx'):
                            file_content = extract_text_from_pptx(io.BytesIO(file.read()))
                        elif file_info.filename.lower().endswith('.docx'):
                            file_content = extract_text_from_docx(io.BytesIO(file.read()))
                        elif file_info.filename.lower().endswith('.pkl'):
                            file_content = extract_text_from_pkl(file)
                        else:
                            file_content = file.read().decode('utf-8')
                        
                        if len(file_content) > 50000:
                            half_content = file_content[:10000]
                            file_content = half_content + "\n\n[Note: The file content was truncated to manage context window size for large information.]"
                        
                        result[relative_path] = file_content
                    except Exception as e:
                        print(f"Skipping file: {relative_path} (unsupported format or error occurred): {e}")
    return result

def count_tokens(text_data):
    # Initialize the tokenizer with GPT-4 or an appropriate model
    tokenizer = GPT2TokenizerFast.from_pretrained('Xenova/gpt-4')
    json_string = json.dumps(text_data)
    tokens = tokenizer.tokenize(json_string)
    return len(tokens)

# Specify the path to your zip file
zip_file_path = '/Users/amirhossain/Desktop/Cornell/NLP/NLP_Project_Four/cs5740-sp24-assignment-4-AmirFone/Hack_lama-main.zip'

# Extract text from the zip file
with zipfile.ZipFile(zip_file_path, 'r') as zip_file:
    text_data = extract_text_from_zip(zip_file, '')

# Save the extracted text data as a JSON file
output_file = 'extracted_text.json'
with open(output_file, 'w') as json_file:
    json.dump(text_data, json_file, indent=4)

# Count tokens
token_count = count_tokens(text_data)

print(f"Text extraction completed. Results saved to {output_file}.")
print(f"Total token count for JSON object: {token_count}")

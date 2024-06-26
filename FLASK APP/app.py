import zipfile
import os
import json
import io
import PyPDF2
from pptx import Presentation
from docx import Document
import pickle
# from transformers import GPT2TokenizerFast
from flask import Flask, request, render_template, jsonify

app = Flask(__name__)

def extract_text_from_pkl(file):
    try:
        data = pickle.load(file)
        if not isinstance(data, str):
            data = str(data)
        return data
    except Exception as e:
        return f"Error occurred during pickle file processing: {e}"

def extract_text_from_pdf(file):
    pdf_reader = PyPDF2.PdfReader(file)
    text = ""
    for page in range(len(pdf_reader.pages)):
        text += pdf_reader.pages[page].extract_text()
    return text

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
    return text

def extract_text_from_docx(file):
    doc = Document(file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_zip(zip_file, file_path):
    result = {}
    skipped_files = []
    for file_info in zip_file.infolist():
        if file_info.filename.startswith(file_path):
            relative_path = file_info.filename[len(file_path):]
            if relative_path and not relative_path.endswith('/'):
                with zip_file.open(file_info) as file:
                    try:
                        file_content = ""
                        if file_info.filename.lower().endswith('.pdf'):
                            file_content = extract_text_from_pdf(io.BytesIO(file.read()))
                        elif file_info.filename.lower().endswith('.pptx'):
                            file_content = extract_text_from_pptx(io.BytesIO(file.read()))
                        elif file_info.filename.lower().endswith('.docx'):
                            file_content = extract_text_from_docx(io.BytesIO(file.read()))
                        elif file_info.filename.lower().endswith('.pkl'):
                            file_content = extract_text_from_pkl(file)
                        else:
                            file_content = file.read().decode('utf-8')

                        if len(file_content) > 50000:
                            half_content = file_content[:10000]
                            file_content = half_content + "\n\n[Note: The file content was truncated to manage context window size for large information.]"

                        result[relative_path] = file_content
                    except Exception as e:
                        skipped_files.append(f"Skipping file: {relative_path} (unsupported format or error occurred): {e}")
    return result, skipped_files

def count_tokens(text_data):
    json_string = json.dumps(text_data)
    total_characters = len(json_string)
    estimated_tokens = total_characters // 4
    return estimated_tokens

@app.route('/', methods=['GET', 'POST'])
def upload_file():
    print('order received')
    if request.method == 'POST':
        uploaded_file = request.files['file']
        if uploaded_file:
            file_extension = os.path.splitext(uploaded_file.filename)[1].lower()
            if file_extension == '.zip':
                with zipfile.ZipFile(io.BytesIO(uploaded_file.read()), 'r') as zip_file:
                    text_data, skipped_files = extract_text_from_zip(zip_file, '')
            else:
                text_data = {uploaded_file.filename: uploaded_file.read().decode('utf-8')}
                skipped_files = []

            token_count = count_tokens(text_data)

            return jsonify({
                'text_data': text_data,
                'token_count': token_count,
                'skipped_files': skipped_files
            })
    print('updating Frontend')
    return render_template('upload.html')

if __name__ == '__main__':
    app.run(debug=True)